from flask import Flask, render_template, request
import os
import google.generativeai as genai  # Replace with your Gemini API library 
import PyPDF2
from pptx import Presentation
from PIL import Image
import requests
import base64
import os
import re

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'  # Folder for temporary file storage
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

genai.configure(api_key=os.getenv("API_KEY")) 

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        if file:
            # 1. File Upload & Processing 
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)
            pages = process_file(file_path)  # Function to extract pages (implement later)

            # 2. Gemini API Calls 
            mermaid_codes = []
            for page_number, page_content in enumerate(pages):
                mermaid_code = generate_mermaid_from_gemini(page_content, page_number)
                mermaid_codes.append(mermaid_code) 

            # 3. Mermaid Diagram Generation (You'll implement these functions later)
            image_paths = generate_mermaid_images(mermaid_codes) 

            # 4. Display and Download 
            return render_template('index.html', image_paths=image_paths)

    return render_template('index.html')

# ... (Implement helper functions below)

def process_file(file_path):
    """
    Handles different file types and extracts page content.

    Args:
        file_path (str): The path to the uploaded file.

    Returns:
        list: A list of strings, where each string represents the content of a page.
    """

    file_extension = os.path.splitext(file_path)[1].lower()
    pages = []

    if file_extension == ".pdf":
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                page_content = page.extract_text()
                pages.append(page_content)

    elif file_extension == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
            pages.append(text)

    elif file_extension in [".jpg", ".jpeg", ".png"]:
        img = Image.open(file_path)
        # For images, we'll assume the entire image is a single "page"
        pages.append(img)  

    elif file_extension == ".txt":
        with open(file_path, 'r') as file:
            pages.append(file.read()) 

    else:
        raise ValueError("Unsupported file type.")

    return pages

config = {
  'temperature': 0,
  'top_k': 20,
  'top_p': 0.9,
  'max_output_tokens': 500
}

safety_settings = [
  {
    "category": "HARM_CATEGORY_HARASSMENT",
    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
  },
  {
    "category": "HARM_CATEGORY_HATE_SPEECH",
    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
  },
  {
    "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
  },
  {
    "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
    "threshold": "BLOCK_MEDIUM_AND_ABOVE"
  }
]

model = genai.GenerativeModel(model_name="gemini-pro",
                              generation_config=config,
                              safety_settings=safety_settings)

def generate_mermaid_from_gemini(page_content, page_number):
    """
    Generates Mermaid code from page content using the Gemini API.

    Args:
        page_content (str or Image): The content of the page (text or image).
        page_number (int): The page number (for reference).

    Returns:
        str: The generated Mermaid code, or a default code if there's an error.
    """

    if isinstance(page_content, str):
        prompt = f"""
        You are a Mermaid diagram expert, specifically skilled in creating mindmaps. 
        Your task is to analyze the following text from page {page_number + 1} and generate Mermaid code that 
        visually represents the key concepts and relationships in the form of a mindmap:
        ```
        {page_content}
        ```
        Make sure the generated Mermaid code is valid and can be rendered correctly by Mermaid.js. 
        Do not include any explanatory text outside of the Mermaid code block. 

        Important Instructions for Mermaid Code Generation:

        * Use the following Mermaid mindmap syntax as a guide: 

        ```
        mindmap
          root((mindmap))
            Origins
              Long history
              ::icon(fa fa-book)
              Popularisation
                British popular psychology author Tony Buzan
            Research
              On effectivness<br/>and features
              On Automatic creation
                Uses
                    Creative techniques
                    Strategic planning
                    Argument mapping
            Tools
              Pen and paper
              Mermaid
        ``` <--- Closing backtick added here
        * Dont repeat the same mindmap as the example make sure to include contents that are actually given
        * Ensure correct indentation and spacing for sub-nodes.
        * Use parentheses `(())` for the root node.
        * Use the correct arrow symbol `-->` for connections where applicable.
        * Include newlines where necessary for proper formatting.
        * Use `::icon(font-awesome icon class)` for adding icons.
        * Use `<br/>` for line breaks within a node.
        """
    else:
        # Handle image content here (if applicable)
        prompt = f"Process the image from page {page_number + 1} to generate Mermaid code that visually represents the key concepts and relationships. Make sure the generated Mermaid code is valid and can be rendered correctly by Mermaid.js. Do not include any explanatory text outside of the Mermaid code block.  "

    try:
        response = model.generate_content(prompt)
        mermaid_code = response.text # This line was likely misaligned 
        # --- Clean up Mermaid code ---
        mermaid_code = mermaid_code.replace('```', '')
        mermaid_code = mermaid_code.strip()
        
        mermaid_code = mermaid_code.replace(' --> ', '-->')  # Remove spaces around arrows
        return mermaid_code
    
    except Exception as e:
        print(f"Error generating Mermaid code for page {page_number + 1}: {e}")
        return "graph TD; A --> B;"

def generate_mermaid_images(mermaid_codes):
    """
    Generates PNG images from Mermaid code snippets using mermaid.ink.
    """

    image_paths = []
    for i, mermaid_code in enumerate(mermaid_codes):
        try:
            # --- Clean up Mermaid code --- 
            mermaid_code = mermaid_code.replace('```', '')
            mermaid_code = mermaid_code.strip()
            mermaid_code = re.sub(r'^mermaid\s+', '', mermaid_code)
            #mermaid_code = "graph LR\n" + mermaid_code
            mermaid_code = mermaid_code.replace(' --> ', '-->')

            # --- Encode Mermaid code ---
            base64_string = base64.b64encode(mermaid_code.encode('utf-8')).decode('utf-8') 

            # --- Construct mermaid.ink URL ---
            url = f"https://mermaid.ink/img/{base64_string}"

            # --- Download the image ---
            response = requests.get(url)
            response.raise_for_status()  

            # --- Save the image ---
            image_name = f"mermaid_diagram_{i}.png"
            image_path = os.path.join("static", image_name)
            print("Saving image to:", image_path)
            with open(image_path, 'wb') as img_file:
                img_file.write(response.content)
            print("Image saved successfully!") 

            image_paths.append(image_name) 

        except Exception as e:
            print(f"Error generating Mermaid image {i}: {e}")

    return image_paths

if __name__ == '__main__':
    app.run(port=5000) 
